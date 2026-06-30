"""Generate BeautyBook Pro investor pitch deck (Uzbek)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / 'BeautyBook_Pro_Investor_Pitch_UZ.pptx'

ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
DARK = RGBColor(0x1E, 0x1B, 0x4B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)


def add_title_slide(prs, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    accent = slide.shapes.add_shape(1, 0, Inches(4.8), prs.slide_width, Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = tf.add_paragraph()
        sub.text = subtitle
        sub.font.size = Pt(22)
        sub.font.color.rgb = RGBColor(0xC4, 0xB5, 0xFD)
        sub.space_before = Pt(16)


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = LIGHT
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(11), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.2), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = MUTED
        p.space_after = Pt(14)


def add_metrics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = 'Bozor va imkoniyat'
    title_box.text_frame.paragraphs[0].font.size = Pt(30)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = DARK

    metrics = [
        ('$2.4 mlrd', 'Global go\'zallik xizmatlari bozori (2025)'),
        ('68%', 'Onlayn bron qilishga o\'tayotgan mijozlar ulushi'),
        ('3x', 'Salonlar uchun qayta bron darajasini oshirish potensiali'),
        ('40%', 'No-show va bo\'sh slotlardan yo\'qotiladigan daromad'),
    ]

    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(5.6)
    height = Inches(2.2)
    gap_x = Inches(0.4)
    gap_y = Inches(0.35)

    for idx, (value, label) in enumerate(metrics):
        row, col = divmod(idx, 2)
        x = left + col * (width + gap_x)
        y = top + row * (height + gap_y)
        card = slide.shapes.add_shape(1, x, y, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = ACCENT

        vbox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.35), width - Inches(0.6), Inches(0.9))
        vp = vbox.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(34)
        vp.font.bold = True
        vp.font.color.rgb = ACCENT

        lbox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.2), width - Inches(0.6), Inches(0.8))
        lp = lbox.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(16)
        lp.font.color.rgb = MUTED
        lp.word_wrap = True


def add_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = 'Investitsiya so\'rovi'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    lines = [
        'Round: Pre-seed — $350,000',
        'Maqsad: O\'zbekiston va Markaziy Osiyoda 500+ salon hamkorligi',
        '12 oy ichida: mobil ilova, to\'lov integratsiyasi, AI slot optimizatsiyasi',
        'Kutilayotgan yillik daromad (ARR): $1.2M (24-oy)',
        'Aloqa: invest@beautybook.uz | +998 90 123 45 67',
    ]
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        para.space_before = Pt(14)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        'BeautyBook Pro',
        'Go\'zallik salonlari uchun aqlli bron va boshqaruv platformasi',
    )

    add_content_slide(prs, 'Muammo', [
        'Salonlar hali ham telefon, Telegram va qog\'oz jadvallar orqali ishlaydi.',
        'Bo\'sh slotlar va no-show tufayli daromadning 30–40% qismi yo\'qoladi.',
        'Mijozlar master va vaqtni tanlashda qiyinchilikka duch keladi.',
        'Salon egalari real vaqt analitikasi va xodimlar samaradorligini ko\'ra olmaydi.',
    ])

    add_content_slide(prs, 'Yechim: BeautyBook Pro', [
        'Mijozlar uchun 4 qadamli onlayn bron (kategoriya → xizmat → master → vaqt).',
        'Masterlar uchun shaxsiy jadval, bo\'sh slotlar va bron tarixi.',
        'Admin panel: bronlar, buyurtmalar, mahsulotlar va hisobotlar.',
        'Do\'kon moduli: kosmetika savdosi va sevimlilar ro\'yxati.',
        'Double-booking himoyasi va JWT autentifikatsiya.',
    ])

    add_content_slide(prs, 'Biznes modeli', [
        'SaaS obuna: salon uchun oylik $29–$99 (xodimlar soniga qarab).',
        'Tranzaksiya komissiyasi: onlayn to\'lovlar bo\'yicha 2.5%.',
        'Premium modullar: SMS eslatmalar, CRM, marketing kampaniyalari.',
        'Marketplace: brend kosmetika mahsulotlaridan marja.',
    ])

    add_metrics_slide(prs)

    add_content_slide(prs, 'Raqobat ustunligi', [
        'Mahalliy til va to\'lov usullariga moslashgan UX (o\'zbek/rus).',
        'Master + mijoz + admin uchun yagona ekotizim.',
        'Real vaqt slot boshqaruvi va conflict detection.',
        'Tez joriy etish: 48 soat ichida salonni tizimga ulash.',
    ])

    add_content_slide(prs, 'Yo\'l xaritasi (12 oy)', [
        'Q1: Toshkent va Samarqand pilot — 50 salon.',
        'Q2: Mobil ilova (iOS/Android) va Payme/Click integratsiyasi.',
        'Q3: AI orqali peak-soat optimizatsiyasi va churn tahlili.',
        'Q4: Qozog\'iston va Qirg\'iziston bozorlariga chiqish.',
    ])

    add_content_slide(prs, 'Jamoa', [
        'CEO — salon tarmoqlari va retail tajribasi (8+ yil).',
        'CTO — Django/DRF va mobil arxitektura mutaxassisi.',
        'COO — operatsion jarayonlar va hamkorliklar bo\'yicha.',
        'Maslahatchilar: go\'zallik industriya ekspertlari va VC mentorlar.',
    ])

    add_ask_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f'Created: {OUTPUT}')


if __name__ == '__main__':
    main()
